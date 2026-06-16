"""
Ajoute des slides manquantes au PowerPoint VillaNova
pour couvrir tous les aspects du projet (20 min de presentation).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from copy import deepcopy
from lxml import etree
import os

# Couleurs du projet
GREEN = RGBColor(0x1B, 0x4D, 0x3E)
TERRACOTTA = RGBColor(0xE7, 0x6F, 0x51)
YELLOW = RGBColor(0xF4, 0xD3, 0x5E)
BEIGE = RGBColor(0xFA, 0xF7, 0xF2)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Emu(18288000)
SLIDE_H = Emu(10287000)

# Marges
LEFT_M = Emu(822960)
TOP_LABEL = Emu(685800)


def add_bg(slide, color):
    """Ajoute un fond de couleur a un slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_section_label(slide, text, color=TERRACOTTA):
    """Ajoute le petit label de section en haut (style 'CODE', 'JAVASCRIPT', etc.)."""
    txBox = slide.shapes.add_textbox(LEFT_M, TOP_LABEL, Emu(8229600), Emu(411480))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    # Espacement entre lettres via texte spaces
    spaced = '   '.join(text.upper())
    run.text = spaced
    run.font.name = 'Consolas'
    run.font.size = Pt(14)
    run.font.color.rgb = color
    run.font.bold = False


def add_title(slide, text, top=Emu(1558290), size=Pt(45), color=DARK):
    """Ajoute le titre principal du slide."""
    txBox = slide.shapes.add_textbox(LEFT_M, top, Emu(16459200), Emu(828675))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = 'Georgia'
    run.font.size = size
    run.font.bold = True
    run.font.color.rgb = color


def add_page_number(slide, num):
    """Ajoute le numero de page en bas a droite."""
    txBox = slide.shapes.add_textbox(Emu(17510760), Emu(9618345), Emu(502920), Emu(348615))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(num)
    run.font.name = 'Consolas'
    run.font.size = Pt(13.5)
    run.font.color.rgb = GRAY


def add_code_block(slide, code_text, left, top, width, height):
    """Ajoute un bloc de code sur fond vert."""
    # Fond vert arrondi
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()
    # Ajustement du rayon
    shape.adjustments[0] = 0.02

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(411480)
    tf.margin_top = Emu(274320)
    tf.margin_right = Emu(411480)
    tf.margin_bottom = Emu(274320)

    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.name = 'Consolas'
        run.font.size = Pt(15)
        run.font.color.rgb = WHITE


def add_text_block(slide, text, left, top, width, height, font_size=Pt(20), color=DARK, bold=False, font_name='Calibri'):
    """Ajoute un bloc de texte."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color


def add_subtitle(slide, text, left, top, width=Emu(6035040)):
    """Sous-titre en gras vert."""
    add_text_block(slide, text, left, top, width, Emu(596265),
                   font_size=Pt(21), color=GREEN, bold=True)


def add_body(slide, text, left, top, width=Emu(6035040)):
    """Texte de corps."""
    add_text_block(slide, text, left, top, width, Emu(600000),
                   font_size=Pt(18), color=GRAY)


def add_multi_text(slide, items, left, top, width, spacing=Emu(800000)):
    """Ajoute plusieurs blocs titre + description empiles."""
    current_top = top
    for title, desc in items:
        add_subtitle(slide, title, left, current_top, width)
        add_body(slide, desc, left, current_top + Emu(420000), width)
        current_top += spacing


def add_file_label(slide, filename, left, top):
    """Label du fichier (ex: 'api.js') en jaune."""
    txBox = slide.shapes.add_textbox(left, top, Emu(3000000), Emu(300000))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = filename
    run.font.name = 'Consolas'
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = YELLOW


# ========================================================
# MAIN
# ========================================================

src = r'c:\Users\Farouk\Documents\LaPlateforme\VillaNova\VilleNova La culture vit ici..pptx'
prs = Presentation(src)
blank_layout = prs.slide_layouts[6]  # Blank layout

# On va ajouter les slides a la fin, puis reordonner.
# Slides actuelles (0-indexed):
#  0: Title
#  1: Contexte
#  2: Ideation
#  3: Croquis
#  4: Wireframe mobile
#  5: Maquette mobile
#  6: Wireframe desktop
#  7: Maquette desktop
#  8: Charte graphique
#  9: HTML semantique
# 10: Architecture SASS
# 11: JS Fetch (page 13 dans PDF)
# 12: Eco-conception (page 12 dans PDF)
# 13: Accessibilite
# 14: Demo
# 15: Merci

new_slides = []

# ---- SLIDE A: Sommaire ----
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, BEIGE)
add_section_label(slide, 'Plan')
add_title(slide, 'Sommaire')

items_left = [
    ('01  Contexte & problématique', DARK),
    ('02  Idéation & recherche', DARK),
    ('03  Maquettes (wireframes & HF)', DARK),
    ('04  Charte graphique', DARK),
    ('05  HTML sémantique', DARK),
    ('06  Architecture SASS', DARK),
]
items_right = [
    ('07  JavaScript modulaire', DARK),
    ('08  API & données dynamiques', DARK),
    ('09  Authentification Supabase', DARK),
    ('10  Accessibilité & éco-conception', DARK),
    ('11  Responsive design', DARK),
    ('12  Démo live', DARK),
]

y = Emu(2800000)
for text, color in items_left:
    add_text_block(slide, text, LEFT_M, y, Emu(7500000), Emu(450000),
                   font_size=Pt(22), color=color, font_name='Calibri')
    y += Emu(550000)

y = Emu(2800000)
for text, color in items_right:
    add_text_block(slide, text, Emu(9500000), y, Emu(7500000), Emu(450000),
                   font_size=Pt(22), color=color, font_name='Calibri')
    y += Emu(550000)

new_slides.append(('sommaire', len(prs.slides._sldIdLst) - 1))


# ---- SLIDE B: Pages du site ----
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, BEIGE)
add_section_label(slide, 'Structure')
add_title(slide, 'Les pages du site')

pages = [
    ('Accueil', 'Hero, mood picker, sélection à la carte,\ncatégories, événements gratuits, engagements,\nnewsletter'),
    ('Tous les événements', 'Grille paginée, filtres par catégorie,\ndate et lieu, bandeau filtres actifs'),
    ('Détail événement', 'Image de couverture, tags, infos pratiques,\ndescription longue, événements liés'),
]
pages2 = [
    ('Catégories', '8 catégories visuelles, résultats dynamiques,\npré-filtrage via URL (?cat=)'),
    ('Carte interactive', 'Carte Leaflet/OSM, marqueurs avec popups,\nauto-zoom sur les résultats'),
    ('Connexion / Inscription', 'Authentification Supabase, email + mot de passe,\nGoogle OAuth, validation formulaire'),
]

x_pos = [LEFT_M, Emu(6400000), Emu(12000000)]
for i, (title, desc) in enumerate(pages):
    add_subtitle(slide, title, x_pos[i], Emu(2700000), Emu(5000000))
    add_body(slide, desc, x_pos[i], Emu(3200000), Emu(5000000))

for i, (title, desc) in enumerate(pages2):
    add_subtitle(slide, title, x_pos[i], Emu(5600000), Emu(5000000))
    add_body(slide, desc, x_pos[i], Emu(6100000), Emu(5000000))

new_slides.append(('pages', len(prs.slides._sldIdLst) - 1))


# ---- SLIDE C: JavaScript - Architecture modulaire ----
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, BEIGE)
add_section_label(slide, 'JavaScript')
add_title(slide, 'Architecture modulaire')

code = """js/
├── api.js            Couche API (fetch OpenAgenda)
├── event-card.js     Factory de cartes + utilitaires partagés
├── home.js           Page d'accueil (hero, mood picker, featured)
├── events-listing.js Page tous les événements (grille, filtres)
├── event-detail.js   Page détail d'un événement
├── categories.js     Page catégories
├── carte.js          Carte interactive Leaflet
├── free-events.js    Section "Gratuit cette semaine"
├── auth.js           Authentification Supabase
└── mobile-menu.js    Menu hamburger mobile"""

add_code_block(slide, code, LEFT_M, Emu(2600000), Emu(10500000), Emu(6200000))

add_multi_text(slide, [
    ('Pattern IIFE', 'Chaque fichier est encapsulé dans une\nIIFE pour éviter la pollution globale.'),
    ('Namespace partagé', 'window.VillaNova expose les utilitaires\ncommuns (createEventCard, formatDate…).'),
    ('Zéro framework', 'Vanilla JS uniquement — DOM API,\nfetch, async/await.'),
], Emu(12000000), Emu(2700000), Emu(5500000))

new_slides.append(('js_archi', len(prs.slides._sldIdLst) - 1))


# ---- SLIDE D: JavaScript - Génération dynamique des cartes ----
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, BEIGE)
add_section_label(slide, 'JavaScript')
add_title(slide, 'Génération dynamique des cartes')

code = """// event-card.js — factory partagée entre toutes les pages

function createEventCard(event, opts) {
  var li = document.createElement('li');
  var card = document.createElement('article');
  card.className = 'event-card';

  // image avec lazy loading
  var img = document.createElement('img');
  img.src = getImageUrl(event, 'thumb');
  img.loading = 'lazy';

  // tags, titre, date, prix...
  var title = document.createElement('h3');
  title.textContent = truncate(getText(event.title), 50);

  li.appendChild(card);
  return li;
}

window.VillaNova.createEventCard = createEventCard;"""

add_code_block(slide, code, LEFT_M, Emu(2600000), Emu(10500000), Emu(6800000))

add_multi_text(slide, [
    ('Une seule factory', 'createEventCard() est utilisée par\n4 pages différentes (accueil, listing,\ncatégories, détail).'),
    ('Utilitaires partagés', 'formatEventDate(), extractPrice(),\ntruncate(), clearChildren()…'),
    ('Lazy loading natif', 'loading="lazy" sur les images\npour réduire le chargement initial.'),
], Emu(12000000), Emu(2700000), Emu(5500000))

new_slides.append(('js_cards', len(prs.slides._sldIdLst) - 1))


# ---- SLIDE E: JavaScript - Filtres & pagination ----
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, BEIGE)
add_section_label(slide, 'JavaScript')
add_title(slide, 'Filtres & pagination')

code = """// events-listing.js — filtres combinés + pagination

function getSearchParams() {
  var urlParams = new URLSearchParams(window.location.search);
  var params = {};

  var category = urlParams.get('category');
  var location = urlParams.get('location');
  var date = urlParams.get('date');

  if (date === 'aujourdhui') {
    params['timings[gte]'] = today;
    params['timings[lte]'] = today;
  } else if (date === 'weekend') {
    params['timings[gte]'] = saturday;
    params['timings[lte]'] = dimanche;
  }
  return params;
}"""

add_code_block(slide, code, LEFT_M, Emu(2600000), Emu(10500000), Emu(6200000))

add_multi_text(slide, [
    ('Filtres par URL', 'Les paramètres ?category=, ?date=,\n?location= sont lus depuis l\'URL\net combinés pour l\'appel API.'),
    ('Pagination offset', 'Bouton "Voir plus" charge 12 événements\nsupplémentaires via offset/limit.'),
    ('Bandeau filtres actifs', 'Affichage visuel des filtres appliqués\navec bouton "Effacer" pour reset.'),
], Emu(12000000), Emu(2700000), Emu(5500000))

new_slides.append(('js_filtres', len(prs.slides._sldIdLst) - 1))


# ---- SLIDE F: JavaScript - Carte Leaflet ----
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, BEIGE)
add_section_label(slide, 'JavaScript')
add_title(slide, 'Carte interactive Leaflet')

code = """// carte.js — intégration Leaflet + OpenStreetMap

var map = L.map('map', { scrollWheelZoom: false })
  .setView([43.2965, 5.3698], 13);

L.tileLayer('https://{s}.tile.openstreetmap.org/...', {
  attribution: '© OpenStreetMap'
}).addTo(map);

// Marqueurs dynamiques depuis l'API
for (var i = 0; i < data.events.length; i++) {
  var ev = data.events[i];
  if (!ev.location.latitude) continue;

  var marker = L.marker([loc.latitude, loc.longitude],
    { icon: eventIcon }).addTo(map);
  marker.bindPopup(buildPopup(ev));
}

// Auto-zoom sur tous les marqueurs
map.fitBounds(L.featureGroup(markers).getBounds());"""

add_code_block(slide, code, LEFT_M, Emu(2600000), Emu(10500000), Emu(6800000))

add_multi_text(slide, [
    ('OpenStreetMap', 'Tuiles gratuites et libres, pas de\nclé API nécessaire.'),
    ('Marqueurs custom', 'DivIcon avec CSS personnalisé\n(classe .map-marker).'),
    ('Popups riches', 'Titre, lieu, date, prix et lien\nvers la page détail.'),
], Emu(12000000), Emu(2700000), Emu(5500000))

new_slides.append(('js_carte', len(prs.slides._sldIdLst) - 1))


# ---- SLIDE G: Authentification Supabase ----
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, BEIGE)
add_section_label(slide, 'Authentification')
add_title(slide, 'Supabase Auth')

code = """// auth.js — connexion email + Google OAuth

var supabase = window.supabase.createClient(
  SUPABASE_URL, SUPABASE_ANON_KEY
);

// Connexion email / mot de passe
async function handleLogin(email, password) {
  var { data, error } = await supabase.auth.signInWithPassword({
    email: email,
    password: password
  });
  if (error) showMsg(error.message, 'error');
  else window.location.href = 'index.html';
}

// Connexion Google OAuth
async function handleGoogleLogin() {
  await supabase.auth.signInWithOAuth({
    provider: 'google',
    options: { redirectTo: window.location.origin }
  });
}"""

add_code_block(slide, code, LEFT_M, Emu(2600000), Emu(10500000), Emu(6800000))

add_multi_text(slide, [
    ('Backend as a Service', 'Supabase fournit l\'auth, la base de\ndonnées et les API — sans serveur\nà gérer.'),
    ('Double méthode', 'Email + mot de passe ET Google OAuth\npour maximiser l\'accessibilité.'),
    ('Gestion de session', 'onAuthStateChange() détecte la session\nactive et adapte le header\n(bouton Connexion → avatar).'),
], Emu(12000000), Emu(2700000), Emu(5500000))

new_slides.append(('supabase', len(prs.slides._sldIdLst) - 1))


# ---- SLIDE H: Responsive Design ----
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, BEIGE)
add_section_label(slide, 'Responsive')
add_title(slide, 'Mobile-first & responsive')

code = """/* _mixins.scss — breakpoints mobile-first */

@mixin tablet  { @media (min-width: 768px)  { @content; } }
@mixin desktop { @media (min-width: 1024px) { @content; } }
@mixin wide    { @media (min-width: 1280px) { @content; } }

/* Utilisation */
.event-card {
  width: 100%;              /* mobile : pleine largeur */

  @include tablet {
    width: calc(50% - 1rem); /* tablette : 2 colonnes */
  }
  @include desktop {
    width: calc(33% - 1rem); /* desktop : 3 colonnes */
  }
}"""

add_code_block(slide, code, LEFT_M, Emu(2600000), Emu(10500000), Emu(5400000))

add_multi_text(slide, [
    ('Mobile-first', 'Le CSS de base cible le mobile.\nLes media queries ajoutent la complexité\nvers le desktop.'),
    ('Bottom nav mobile', 'Barre de navigation fixe en bas\nde l\'écran sur mobile avec 5 onglets\n(Accueil, Explorer, Carte, Favoris, Compte).'),
    ('Menu hamburger', 'Menu latéral avec overlay, gestion\ndu focus-trap et fermeture via Échap.'),
    ('Grilles adaptatives', 'CSS Grid et Flexbox pour des layouts\nfluides selon la taille d\'écran.'),
], Emu(12000000), Emu(2600000), Emu(5500000), Emu(700000))

new_slides.append(('responsive', len(prs.slides._sldIdLst) - 1))


# ---- SLIDE I: API OpenAgenda (plus détaillé) ----
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, BEIGE)
add_section_label(slide, 'API')
add_title(slide, 'Intégration OpenAgenda')

code = """// api.js — couche d'abstraction API

var API_BASE = 'https://api.openagenda.com/v2/agendas/69319016';
var API_KEY  = '...';

function buildUrl(path, params) {
  var url = new URL(API_BASE + path);
  params.key = API_KEY;
  for (var k in params) url.searchParams.set(k, params[k]);
  return url.toString();
}

async function fetchEvents(params) {
  var url = buildUrl('/events', params);
  var res = await fetch(url);
  if (!res.ok) throw new Error('API ' + res.status);
  var json = await res.json();
  return { events: json.events, total: json.total };
}"""

add_code_block(slide, code, LEFT_M, Emu(2600000), Emu(10500000), Emu(6200000))

add_multi_text(slide, [
    ('API RESTful', 'OpenAgenda expose les événements\nde la ville via une API JSON publique.'),
    ('Paramètres de filtrage', 'search, timings[gte], timings[lte],\nlimit, offset — combinables.'),
    ('Données multilingues', 'Titres et descriptions en fr/en,\nextraits via getText(field).'),
], Emu(12000000), Emu(2700000), Emu(5500000))

new_slides.append(('api', len(prs.slides._sldIdLst) - 1))


# ---- SLIDE J: Bilan technique ----
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, BEIGE)
add_section_label(slide, 'Bilan')
add_title(slide, 'Bilan technique')

# Colonne gauche
stats = [
    ('6', 'pages HTML'),
    ('10', 'fichiers JavaScript'),
    ('12', 'fichiers SASS'),
    ('0', 'framework JS'),
]

x = LEFT_M
for val, label in stats:
    txBox = slide.shapes.add_textbox(x, Emu(2800000), Emu(3500000), Emu(1200000))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = val
    run.font.name = 'Georgia'
    run.font.size = Pt(60)
    run.font.bold = True
    run.font.color.rgb = GREEN

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = label
    run2.font.name = 'Calibri'
    run2.font.size = Pt(18)
    run2.font.color.rgb = GRAY
    x += Emu(4000000)

# Technologies
add_subtitle(slide, 'Technologies utilisées', LEFT_M, Emu(5200000), Emu(16000000))

techs = [
    ('HTML5 / CSS3 / SASS', 'Structure, styles, préprocesseur'),
    ('JavaScript ES6+', 'Vanilla JS, async/await, DOM API'),
    ('Supabase', 'Auth (email + Google OAuth)'),
    ('Leaflet + OpenStreetMap', 'Carte interactive'),
    ('OpenAgenda API', 'Données événementielles'),
]

x = LEFT_M
for title, desc in techs:
    add_text_block(slide, title, x, Emu(5800000), Emu(3200000), Emu(400000),
                   font_size=Pt(18), color=DARK, bold=True)
    add_text_block(slide, desc, x, Emu(6200000), Emu(3200000), Emu(400000),
                   font_size=Pt(15), color=GRAY)
    x += Emu(3400000)

new_slides.append(('bilan', len(prs.slides._sldIdLst) - 1))


# ========================================================
# REORDONNER LES SLIDES
# ========================================================

# Ordre voulu :
#  0:  Title (existant 0)
#  1:  Sommaire (nouveau A)
#  2:  Contexte (existant 1)
#  3:  Ideation (existant 2)
#  4:  Pages du site (nouveau B)
#  5:  Croquis (existant 3)
#  6:  Wireframe mobile (existant 4)
#  7:  Maquette mobile (existant 5)
#  8:  Wireframe desktop (existant 6)
#  9:  Maquette desktop (existant 7)
# 10:  Charte graphique (existant 8)
# 11:  HTML semantique (existant 9)
# 12:  SASS (existant 10)
# 13:  JS Archi modulaire (nouveau C)
# 14:  JS Fetch/API (existant 11) => remplacé par nouveau I (API OpenAgenda)
# 15:  JS Cartes dynamiques (nouveau D)
# 16:  JS Filtres (nouveau E)
# 17:  JS Carte Leaflet (nouveau F)
# 18:  Supabase (nouveau G)
# 19:  Eco-conception (existant 12)
# 20:  Accessibilite (existant 13)
# 21:  Responsive (nouveau H)
# 22:  Bilan technique (nouveau J)
# 23:  Demo (existant 14)
# 24:  Merci (existant 15)

# Les slides existantes sont indices 0-15 dans l'original
# Les nouvelles sont 16-24 (A=16, B=17, C=18, D=19, E=20, F=21, G=22, H=23, I=24, J=25)

desired_order = [
    0,   # Title
    16,  # Sommaire (A)
    1,   # Contexte
    2,   # Ideation
    17,  # Pages du site (B)
    3,   # Croquis
    4,   # Wireframe mobile
    5,   # Maquette mobile
    6,   # Wireframe desktop
    7,   # Maquette desktop
    8,   # Charte graphique
    9,   # HTML semantique
    10,  # SASS
    18,  # JS Archi modulaire (C)
    24,  # API OpenAgenda (I) - remplace l'ancien JS Fetch
    19,  # JS Cartes dynamiques (D)
    20,  # JS Filtres (E)
    21,  # JS Carte Leaflet (F)
    22,  # Supabase (G)
    12,  # Eco-conception
    13,  # Accessibilite
    23,  # Responsive (H)
    25,  # Bilan technique (J)
    14,  # Demo
    15,  # Merci
]

# Reorder slides using XML manipulation
sldIdLst = prs._element.find(
    '{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst'
)
sldIds = list(sldIdLst)

# Build a mapping of current position to sldId element
reordered = []
for idx in desired_order:
    if idx < len(sldIds):
        reordered.append(sldIds[idx])

# Clear and re-add in order
for sldId in list(sldIdLst):
    sldIdLst.remove(sldId)

for sldId in reordered:
    sldIdLst.append(sldId)


# ========================================================
# SAVE
# ========================================================

out = r'c:\Users\Farouk\Documents\LaPlateforme\VillaNova\VillaNova_presentation_complete.pptx'
prs.save(out)
print(f'Presentation sauvegardee: {out}')
print(f'Nombre total de slides: {len(prs.slides)}')
